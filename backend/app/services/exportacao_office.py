"""Exportação universal em formatos Office editáveis.

O conteúdo chega já resolvido a partir dos registros canônicos publicados. Este
módulo muda apenas a diagramação: não recebe texto clínico do navegador, não
resume e não cria recomendações. PPTX e DOCX preservam títulos, seções, listas,
referências e proveniência como texto nativo editável.
"""
from __future__ import annotations

from datetime import datetime, timezone
from io import BytesIO
from typing import Any

from docx import Document as WordDocument
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DocxInches
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.services.pdf.marca import LOGO, logo_disponivel
from app.services.professional_profile import (
    document_identity,
    professional_name,
    workplace_lines,
)

from .apresentacao import _fragmentar
from .exportacao_conteudo import (
    ROTULOS_TIPO,
    ConteudoExportavel,
    _registro_profissional,
    _sem_markdown,
)


NAVY = RGBColor(0x0B, 0x2E, 0x45)
TINTA = RGBColor(0x26, 0x33, 0x3B)
NEUTRO = RGBColor(0x55, 0x66, 0x6F)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
LARGURA = Inches(13.333)
ALTURA = Inches(7.5)
MARGEM = Inches(0.65)
MAX_BLOCOS_SLIDE = 6


def _titulo_exportacao(itens: list[ConteudoExportavel], titulo: str | None) -> str:
    if titulo and titulo.strip():
        return titulo.strip()[:180]
    if len(itens) == 1:
        return itens[0].titulo
    return "Seleção de conteúdo CorVIA"


def _identificacao(user: Any, incluir: bool) -> tuple[str, list[str]]:
    if not incluir:
        return "CorVIA", []
    identidade = document_identity(user)
    nome = professional_name(identidade) or "Assinante CorVIA"
    linhas = [x for x in [_registro_profissional(identidade), identidade.get("specialty")] if x]
    linhas.extend(workplace_lines(identidade))
    return nome, [str(x) for x in linhas if x]


def _slide_vazio(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _caixa(slide, esquerda, topo, largura, altura):
    shape = slide.shapes.add_textbox(esquerda, topo, largura, altura)
    shape.text_frame.word_wrap = True
    return shape.text_frame


def _rodape_slide(slide, texto: str) -> None:
    if logo_disponivel():
        try:
            slide.shapes.add_picture(str(LOGO), MARGEM, ALTURA - Inches(0.42), width=Inches(1.3))
        except (OSError, ValueError):
            pass
    tf = _caixa(slide, Inches(2.05), ALTURA - Inches(0.42), LARGURA - Inches(2.65), Inches(0.28))
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(9)
    p.font.color.rgb = NEUTRO


def _capa_pptx(prs: Presentation, titulo: str, subtitulo: str, nome: str, detalhes: list[str]) -> None:
    slide = _slide_vazio(prs)
    fundo = slide.shapes.add_shape(1, 0, 0, LARGURA, ALTURA)
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = NAVY
    fundo.line.fill.background()
    if logo_disponivel():
        try:
            slide.shapes.add_picture(str(LOGO), Inches(0.85), Inches(0.55), width=Inches(2.45))
        except (OSError, ValueError):
            pass
    tf = _caixa(slide, Inches(0.85), Inches(2.05), LARGURA - Inches(1.7), Inches(2.2))
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    tf2 = _caixa(slide, Inches(0.85), Inches(4.25), LARGURA - Inches(1.7), Inches(0.9))
    p2 = tf2.paragraphs[0]
    p2.text = subtitulo
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xCF, 0xDD, 0xE4)
    tf3 = _caixa(slide, Inches(0.85), ALTURA - Inches(1.05), LARGURA - Inches(1.7), Inches(0.7))
    p3 = tf3.paragraphs[0]
    p3.text = " · ".join([nome, *detalhes, "CorVIA — Clinical OS"])
    p3.font.size = Pt(11)
    p3.font.color.rgb = RGBColor(0xA9, 0xBD, 0xC8)


def _slide_conteudo(prs: Presentation, titulo: str, blocos: list[tuple[str, str]], rodape: str) -> None:
    slide = _slide_vazio(prs)
    faixa = slide.shapes.add_shape(1, 0, 0, LARGURA, Inches(0.95))
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = NAVY
    faixa.line.fill.background()
    tf_titulo = faixa.text_frame
    tf_titulo.margin_left = MARGEM
    tf_titulo.margin_top = Inches(0.14)
    p_titulo = tf_titulo.paragraphs[0]
    p_titulo.text = titulo
    p_titulo.font.size = Pt(27)
    p_titulo.font.bold = True
    p_titulo.font.color.rgb = BRANCO

    tf = _caixa(slide, MARGEM, Inches(1.2), LARGURA - 2 * MARGEM, ALTURA - Inches(1.9))
    for indice, (tipo, texto) in enumerate(blocos):
        p = tf.paragraphs[0] if indice == 0 else tf.add_paragraph()
        p.text = texto
        p.level = 0
        p.font.size = Pt(19 if tipo == "item" else 18)
        p.font.color.rgb = TINTA
        p.space_after = Pt(9)
        if tipo == "item":
            p.text = f"•  {texto}"
    _rodape_slide(slide, rodape)


def _blocos_secao(paragrafos: list[str], itens: list[str], destaque: str | None) -> list[tuple[str, str]]:
    blocos: list[tuple[str, str]] = []
    if destaque:
        blocos.extend(("paragrafo", parte) for parte in _fragmentar(_sem_markdown(destaque), 260))
    for paragrafo in paragrafos:
        blocos.extend(("paragrafo", parte) for parte in _fragmentar(_sem_markdown(paragrafo), 260))
    for item in itens:
        texto = _sem_markdown(item)
        if texto:
            blocos.extend(("item", parte) for parte in _fragmentar(texto, 230))
    return blocos


def gerar_pptx(
    itens: list[ConteudoExportavel],
    *,
    user: Any,
    incluir_dados_assinante: bool,
    titulo: str | None = None,
) -> bytes:
    if not itens:
        raise ValueError("Nenhum conteúdo elegível para exportação.")
    titulo_final = _titulo_exportacao(itens, titulo)
    nome, detalhes = _identificacao(user, incluir_dados_assinante)
    prs = Presentation()
    prs.slide_width = LARGURA
    prs.slide_height = ALTURA
    subtitulo = itens[0].tema or ROTULOS_TIPO[itens[0].tipo] if len(itens) == 1 else f"{len(itens)} conteúdos selecionados"
    _capa_pptx(prs, titulo_final, subtitulo, nome, detalhes)
    rodape = f"CorVIA — Clinical OS · exportado em {datetime.now(timezone.utc).strftime('%d/%m/%Y')}"

    _slide_conteudo(prs, "Proveniência", [
        ("paragrafo", "Este arquivo reproduz conteúdo publicado no CorVIA na data da exportação."),
        ("paragrafo", "Referências e proveniência acompanham cada seção quando disponíveis."),
        ("paragrafo", "A exportação não cria nem atualiza recomendações clínicas."),
    ], rodape)

    for numero_item, item in enumerate(itens, start=1):
        prefixo = f"{numero_item}. " if len(itens) > 1 else ""
        metadados = " · ".join(x for x in [ROTULOS_TIPO[item.tipo], item.tema, item.subtitulo] if x)
        _slide_conteudo(prs, f"{prefixo}{item.titulo}", [("paragrafo", metadados or ROTULOS_TIPO[item.tipo])], rodape)
        for secao in item.secoes:
            blocos = _blocos_secao(secao.paragrafos, secao.itens, secao.destaque)
            if not blocos:
                continue
            paginas = [blocos[i:i + MAX_BLOCOS_SLIDE] for i in range(0, len(blocos), MAX_BLOCOS_SLIDE)]
            for pagina, bloco in enumerate(paginas, start=1):
                rotulo = secao.titulo if len(paginas) == 1 else f"{secao.titulo} ({pagina}/{len(paginas)})"
                _slide_conteudo(prs, rotulo, bloco, rodape)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _definir_borda_inferior(paragrafo, cor: str = "1C7293", tamanho: str = "12") -> None:
    p_pr = paragrafo._p.get_or_add_pPr()
    p_bdr = p_pr.find(qn("w:pBdr"))
    if p_bdr is None:
        p_bdr = OxmlElement("w:pBdr")
        p_pr.append(p_bdr)
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), tamanho)
    bottom.set(qn("w:space"), "5")
    bottom.set(qn("w:color"), cor)
    p_bdr.append(bottom)


def gerar_docx(
    itens: list[ConteudoExportavel],
    *,
    user: Any,
    incluir_dados_assinante: bool,
    titulo: str | None = None,
) -> bytes:
    if not itens:
        raise ValueError("Nenhum conteúdo elegível para exportação.")
    titulo_final = _titulo_exportacao(itens, titulo)
    nome, detalhes = _identificacao(user, incluir_dados_assinante)
    documento = WordDocument()
    secao = documento.sections[0]
    secao.page_width = DocxInches(8.27)
    secao.page_height = DocxInches(11.69)
    secao.top_margin = DocxInches(0.8)
    secao.bottom_margin = DocxInches(0.75)
    secao.left_margin = DocxInches(0.85)
    secao.right_margin = DocxInches(0.85)

    normal = documento.styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = DocxPt(10.5)
    normal.paragraph_format.space_after = DocxPt(6)
    for estilo, tamanho, cor in (("Title", 28, "0B2E45"), ("Heading 1", 19, "0B2E45"), ("Heading 2", 14, "1C7293")):
        style = documento.styles[estilo]
        style.font.name = "Aptos Display" if estilo != "Normal" else "Aptos"
        style.font.size = DocxPt(tamanho)
        style.font.color.rgb = DocxRGBColor.from_string(cor)

    if logo_disponivel():
        try:
            documento.add_picture(str(LOGO), width=DocxInches(1.7))
        except (OSError, ValueError):
            pass
    titulo_p = documento.add_paragraph(style="Title")
    titulo_p.add_run(titulo_final)
    _definir_borda_inferior(titulo_p)
    subtitulo = itens[0].tema or ROTULOS_TIPO[itens[0].tipo] if len(itens) == 1 else f"{len(itens)} conteúdos selecionados"
    p_sub = documento.add_paragraph()
    run_sub = p_sub.add_run(subtitulo)
    run_sub.bold = True
    run_sub.font.size = DocxPt(12)
    run_sub.font.color.rgb = DocxRGBColor.from_string("55666F")
    p_ident = documento.add_paragraph()
    p_ident.add_run(" · ".join([nome, *detalhes])).italic = True

    documento.add_heading("Proveniência", level=1)
    documento.add_paragraph(
        "Este arquivo reproduz conteúdo publicado no CorVIA na data da exportação. "
        "Referências e proveniência acompanham cada seção quando disponíveis. "
        "A exportação não cria nem atualiza recomendações clínicas."
    )

    for indice, item in enumerate(itens, start=1):
        if indice > 1:
            documento.add_page_break()
        titulo_item = f"{indice}. {item.titulo}" if len(itens) > 1 else item.titulo
        documento.add_heading(titulo_item, level=1)
        metadados = " · ".join(x for x in [ROTULOS_TIPO[item.tipo], item.tema, item.subtitulo] if x)
        if metadados:
            p_meta = documento.add_paragraph()
            run_meta = p_meta.add_run(metadados)
            run_meta.bold = True
            run_meta.font.color.rgb = DocxRGBColor.from_string("55666F")
        for secao_item in item.secoes:
            documento.add_heading(secao_item.titulo, level=2)
            if secao_item.destaque:
                p = documento.add_paragraph()
                r = p.add_run(_sem_markdown(secao_item.destaque))
                r.bold = True
                r.font.color.rgb = DocxRGBColor.from_string("0B2E45")
            for paragrafo in secao_item.paragrafos:
                texto = _sem_markdown(paragrafo)
                if texto:
                    documento.add_paragraph(texto)
            for item_texto in secao_item.itens:
                texto = _sem_markdown(item_texto)
                if texto:
                    documento.add_paragraph(texto, style="List Bullet")

    for secao_doc in documento.sections:
        footer = secao_doc.footer.paragraphs[0]
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = footer.add_run(
            f"CorVIA — Clinical OS · exportado em {datetime.now(timezone.utc).strftime('%d/%m/%Y')}"
        )
        run.font.size = DocxPt(8)
        run.font.color.rgb = DocxRGBColor.from_string("55666F")

    documento.core_properties.title = titulo_final
    documento.core_properties.author = nome
    documento.core_properties.subject = "Conteúdo clínico exportado do CorVIA"
    buffer = BytesIO()
    documento.save(buffer)
    return buffer.getvalue()
