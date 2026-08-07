"""Modo apresentação: exporta o mesmo conteúdo do PDF (`apresentacao.py`) em
PowerPoint (.pptx) nativo e editável — pedido do Rafael em 07/08/2026, para
o assinante poder ajustar os slides antes de projetar (reordenar, cortar,
adicionar algo próprio), o que um PDF não permite.

Reaproveita deliberadamente a mesma extração de conteúdo do PDF
(`_secoes`, `_limpar`, `_fragmentar`, `MARCADORES_POR_PAGINA`, a árvore de
decisão via `arv`) — o texto que sai aqui é o mesmo, verificado, que já sai
no PDF; só o formato de arquivo muda. Nenhum conteúdo é gerado ou reescrito
para esta exportação.

A árvore de decisão vira lista com indentação (não um desenho) — um .pptx
serve para o assinante editar, e uma lista de texto é o que dá para editar de
verdade; o desenho fiel ao Mermaid já existe no PDF e no fluxograma da tela.
"""
from __future__ import annotations

from datetime import datetime, timezone
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.models.content import Document
from app.services.professional_profile import professional_name, workplace_lines

from .apresentacao import MARCADORES_POR_PAGINA, _fragmentar, _limpar, _secoes
from .pdf import arvore as arv

NAVY = RGBColor(0x0B, 0x2E, 0x45)
TEAL = RGBColor(0x1C, 0x72, 0x93)
VERMELHO = RGBColor(0xD5, 0x00, 0x1D)
TINTA = RGBColor(0x26, 0x33, 0x3B)
NEUTRO = RGBColor(0x55, 0x66, 0x6F)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)

LARGURA = Inches(13.333)
ALTURA = Inches(7.5)
MARGEM = Inches(0.6)


def _slide_em_branco(prs: Presentation):
    layout_em_branco = prs.slide_layouts[6]  # layout "Blank" do template padrão
    return prs.slides.add_slide(layout_em_branco)


def _caixa_texto(slide, esquerda, topo, largura, altura):
    caixa = slide.shapes.add_textbox(esquerda, topo, largura, altura)
    tf = caixa.text_frame
    tf.word_wrap = True
    return tf


def _rodape(slide, texto: str) -> None:
    tf = _caixa_texto(slide, MARGEM, ALTURA - Inches(0.45), LARGURA - 2 * MARGEM, Inches(0.35))
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(10)
    p.font.color.rgb = NEUTRO


def _titulo_slide(slide, titulo: str) -> None:
    faixa = slide.shapes.add_shape(1, 0, 0, LARGURA, Inches(0.95))  # MSO_SHAPE.RECTANGLE = 1
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = NAVY
    faixa.line.fill.background()
    faixa.shadow.inherit = False
    tf = faixa.text_frame
    tf.word_wrap = True
    tf.margin_left = MARGEM
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRANCO


def _corpo_com_marcadores(slide, blocos: list[tuple[str, str]]) -> None:
    tf = _caixa_texto(slide, MARGEM, Inches(1.25), LARGURA - 2 * MARGEM, ALTURA - Inches(2.0))
    primeiro = True
    for tipo, texto in blocos:
        p = tf.paragraphs[0] if primeiro else tf.add_paragraph()
        primeiro = False
        p.text = f"•  {texto}" if tipo == "item" else texto
        p.font.size = Pt(20 if tipo == "item" else 18)
        p.font.color.rgb = TINTA
        p.space_after = Pt(10)


def _capa(prs: Presentation, titulo: str, resumo: str, nome: str, registro: str) -> None:
    slide = _slide_em_branco(prs)
    fundo = slide.shapes.add_shape(1, 0, 0, LARGURA, ALTURA)
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = NAVY
    fundo.line.fill.background()
    fundo.shadow.inherit = False

    tf = _caixa_texto(slide, Inches(0.9), Inches(2.2), LARGURA - Inches(1.8), Inches(2.2))
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    if resumo:
        tf2 = _caixa_texto(slide, Inches(0.9), Inches(4.1), LARGURA - Inches(1.8), Inches(1.4))
        p2 = tf2.paragraphs[0]
        p2.text = _limpar(resumo)
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xCF, 0xDD, 0xE4)

    rodape = " · ".join(x for x in (nome, registro) if x) or "Corvia"
    tf3 = _caixa_texto(slide, Inches(0.9), ALTURA - Inches(0.9), LARGURA - Inches(1.8), Inches(0.5))
    p3 = tf3.paragraphs[0]
    p3.text = f"{rodape}  ·  Corvia · corvia.med.br"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(0x9D, 0xB4, 0xC2)


def gerar(doc: Document, medico: dict, anotacao: str = "") -> bytes:
    """Mesma composição do PDF (`apresentacao.gerar`), em .pptx. Ver o
    docstring do módulo — nenhum texto novo é produzido aqui."""
    nome = professional_name(medico)
    conselho = " ".join(
        parte for parte in (medico.get("council_name"), medico.get("council_number")) if parte
    )
    if medico.get("council_state"):
        conselho = f"{conselho}/{medico['council_state']}".strip("/")
    registro = " · ".join(
        parte for parte in (
            conselho, f"RQE {medico['rqe']}" if medico.get("rqe") else "", *workplace_lines(medico),
        ) if parte
    )

    prs = Presentation()
    prs.slide_width = LARGURA
    prs.slide_height = ALTURA

    _capa(prs, doc.title, doc.summary or "", nome or "Corvia", registro)

    rodape_txt = f"Corvia · {doc.theme} · corvia.med.br"
    paginas_conteudo = 0

    def abrir_slide(titulo: str):
        nonlocal paginas_conteudo
        slide = _slide_em_branco(prs)
        _titulo_slide(slide, titulo)
        _rodape(slide, rodape_txt)
        paginas_conteudo += 1
        return slide

    fonte = arv.extrair_mermaid(doc.body_md or "")
    if fonte:
        raiz = arv.montar(fonte)
        slide = abrir_slide("Árvore de decisão")
        if raiz is None:
            _corpo_com_marcadores(slide, [("paragrafo", "O diagrama não pôde ser convertido em lista com fidelidade. Consulte o fluxograma original na Corvia.")])
        else:
            _corpo_com_marcadores(slide, [("item", linha) for linha in arv.em_lista(raiz)])

    secoes = _secoes(doc.body_md or "")
    for titulo, corpo in secoes:
        if not corpo:
            continue
        blocos_pag = [
            corpo[indice:indice + MARCADORES_POR_PAGINA]
            for indice in range(0, len(corpo), MARCADORES_POR_PAGINA)
        ]
        for numero, bloco in enumerate(blocos_pag, start=1):
            rotulo = titulo or doc.title
            if len(blocos_pag) > 1:
                rotulo = f"{rotulo} ({numero}/{len(blocos_pag)})"
            slide = abrir_slide(rotulo)
            _corpo_com_marcadores(slide, bloco)

    if paginas_conteudo == 0:
        slide = abrir_slide("Resumo")
        resumo = _limpar(doc.summary or "")
        if resumo:
            _corpo_com_marcadores(slide, [("paragrafo", f) for f in _fragmentar(resumo)])
        else:
            _corpo_com_marcadores(slide, [(
                "paragrafo",
                "O documento publicado não contém corpo textual projetável. Consulte o "
                "conteúdo original e suas referências na Biblioteca científica.",
            )])

    if anotacao.strip():
        slide = abrir_slide("Anotação do apresentador")
        linhas = [("paragrafo", _limpar(p)) for p in anotacao.strip().splitlines() if p.strip()]
        _corpo_com_marcadores(slide, linhas)
        tf = _caixa_texto(slide, MARGEM, ALTURA - Inches(0.85), LARGURA - 2 * MARGEM, Inches(0.35))
        p = tf.paragraphs[0]
        p.text = "Anotação de quem exportou — não faz parte do conteúdo verificado da plataforma."
        p.font.italic = True
        p.font.size = Pt(10)
        p.font.color.rgb = NEUTRO

    slide = abrir_slide("Procedência")
    blocos: list[tuple[str, str]] = [
        ("paragrafo", f"Documento: {doc.title}"),
        ("paragrafo", f"Tema: {doc.theme} · nível de fonte: {doc.source_tier} · revisão: {doc.review_status}"),
    ]
    if doc.source_refs:
        blocos.extend(("item", _limpar(r)) for r in doc.source_refs[:8])
        if len(doc.source_refs) > 8:
            blocos.append(("paragrafo", f"e mais {len(doc.source_refs) - 8} referência(s) no documento original."))
    if doc.gaps:
        blocos.append(("paragrafo", f"⚠ Este documento declara {len(doc.gaps)} ponto(s) pendente(s) de verificação humana."))
    hoje = datetime.now(timezone.utc).astimezone().strftime("%d/%m/%Y")
    blocos.append(("paragrafo", f"Exportado da Corvia em {hoje}" + (f" por {nome}." if nome else ".")))
    _corpo_com_marcadores(slide, blocos)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
