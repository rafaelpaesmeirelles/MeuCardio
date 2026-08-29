from __future__ import annotations

import hashlib
import io
import json
import re
import unicodedata
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import fitz
import httpx
from docx import Document as DocxDocument
from pptx import Presentation
from sqlalchemy.orm import Session

from app.core.config import settings
from app.models.content import Document
from app.models.scientific_user_document import ScientificUserDocument
from app.services import cofre
from app.services.knowledge_graph import backfill_mesmo_tema

RESPONSES_URL = "https://api.openai.com/v1/responses"
MAX_EXTRACTED_CHARS = 300_000
ANALYSIS_CONTEXT_CHARS = 120_000
TRANSLATION_CHUNK_CHARS = 12_000

ANALYSIS_SCHEMA = {
    "type": "object",
    "additionalProperties": False,
    "properties": {
        "title": {"type": "string"},
        "document_type": {
            "type": "string",
            "enum": ["diretriz", "consenso", "ensaio_clinico", "estudo_observacional", "revisao_sistematica", "meta_analise", "evidencia", "artigo", "outro"],
        },
        "language": {"type": "string"},
        "doi": {"type": ["string", "null"]},
        "source_url": {"type": ["string", "null"]},
        "summary_pt": {"type": "string"},
        "methodology_pt": {"type": "string"},
        "population_pt": {"type": "string"},
        "interventions_pt": {"type": "string"},
        "outcomes_pt": {"type": "string"},
        "results_pt": {"type": "string"},
        "limitations_pt": {"type": "array", "items": {"type": "string"}},
        "clinical_implications_pt": {"type": "array", "items": {"type": "string"}},
        "key_points_pt": {"type": "array", "items": {"type": "string"}},
        "topics": {"type": "array", "items": {"type": "string"}},
        "evidence_strength": {"type": "string", "enum": ["alta", "moderada", "baixa", "muito_baixa", "nao_aplicavel", "nao_determinada"]},
        "needs_translation": {"type": "boolean"},
        "adds_to_corvia": {"type": "boolean"},
        "incorporation_reason_pt": {"type": "string"},
    },
    "required": [
        "title", "document_type", "language", "doi", "source_url", "summary_pt",
        "methodology_pt", "population_pt", "interventions_pt", "outcomes_pt",
        "results_pt", "limitations_pt", "clinical_implications_pt", "key_points_pt",
        "topics", "evidence_strength", "needs_translation", "adds_to_corvia",
        "incorporation_reason_pt",
    ],
}


def private_root() -> Path:
    return Path(settings.exames_dir) / "scientific-user-library"


def sha256_bytes(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


def _pdf_text(content: bytes) -> str:
    doc = fitz.open(stream=content, filetype="pdf")
    try:
        return "\n\n".join(page.get_text("text") for page in doc)
    finally:
        doc.close()


def _docx_text(content: bytes) -> str:
    doc = DocxDocument(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _pptx_text(content: bytes) -> str:
    presentation = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text)
    return "\n\n".join(parts)


def extract_text(content: bytes, media_type: str) -> str:
    if media_type == "application/pdf":
        text = _pdf_text(content)
    elif media_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        text = _docx_text(content)
    elif media_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = _pptx_text(content)
    elif media_type in {"text/plain", "text/csv"}:
        text = content.decode("utf-8-sig")
    else:
        raise ValueError("Formato científico ainda não possui extração textual segura.")
    text = text.replace("\x00", "").strip()
    if not text:
        raise ValueError("O arquivo não contém texto extraível. PDFs somente-imagem exigem análise multimodal específica.")
    return text[:MAX_EXTRACTED_CHARS]


def _model() -> str:
    return (settings.ai_cardiovascular_exam_model or settings.openai_model or "gpt-5.6").strip()


def _response_text(payload: dict) -> str:
    parts: list[str] = []
    for output in payload.get("output") or []:
        if output.get("type") != "message":
            continue
        for item in output.get("content") or []:
            if item.get("type") == "output_text":
                parts.append(str(item.get("text") or ""))
    return "".join(parts).strip()


def _post_response(request: dict, timeout: float = 220.0) -> dict:
    if not settings.ai_enabled or settings.ai_provider != "openai" or not settings.openai_api_key.strip():
        raise RuntimeError("Provedor de IA clínica não configurado.")
    with httpx.Client(timeout=httpx.Timeout(timeout, connect=20.0)) as client:
        response = client.post(
            RESPONSES_URL,
            headers={
                "Authorization": f"Bearer {settings.openai_api_key}",
                "Content-Type": "application/json",
            },
            json=request,
        )
    response.raise_for_status()
    payload = response.json()
    if payload.get("status") != "completed" or payload.get("error"):
        raise ValueError("O provedor não concluiu a análise científica.")
    return payload


def analyze_text(text: str) -> dict[str, Any]:
    prompt = text[:ANALYSIS_CONTEXT_CHARS]
    request = {
        "model": _model(),
        "store": False,
        "max_output_tokens": max(4096, settings.ai_max_output_tokens),
        "instructions": (
            "Você é o analisador científico do CorVIA. Analise exclusivamente o documento fornecido. "
            "Não invente dados, doses, classes de recomendação, níveis de evidência ou resultados. "
            "Diferencie claramente resultados do estudo, limitações e implicações clínicas. Produza a "
            "síntese em português do Brasil. DOI e URL só podem ser informados quando constarem no próprio "
            "documento. Identifique se o documento acrescenta conhecimento útil ao acervo CorVIA, mas não "
            "autorize incorporação: isso depende de consentimento explícito do usuário."
        ),
        "input": [{"role": "user", "content": [{"type": "input_text", "text": prompt}]}],
        "text": {"format": {"type": "json_schema", "name": "corvia_scientific_document_analysis", "strict": True, "schema": ANALYSIS_SCHEMA}},
    }
    payload = _post_response(request)
    raw = _response_text(payload)
    try:
        analysis = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ValueError("O provedor devolveu análise científica inválida.") from exc
    analysis["analyzed_at"] = datetime.now(timezone.utc).isoformat()
    analysis["model"] = _model()
    return analysis


def _translate_chunk(chunk: str) -> str:
    request = {
        "model": _model(),
        "store": False,
        "max_output_tokens": max(4096, settings.ai_max_output_tokens),
        "instructions": (
            "Traduza fielmente para português do Brasil o trecho científico fornecido para uso privado do "
            "usuário que enviou o arquivo. Preserve títulos, subtítulos, números, unidades, tabelas em texto, "
            "abreviações e referências. Não resuma, não interprete e não acrescente conteúdo. Retorne somente a tradução."
        ),
        "input": [{"role": "user", "content": [{"type": "input_text", "text": chunk}]}],
    }
    return _response_text(_post_response(request))


def translate_full_text(text: str) -> str:
    chunks = [text[i:i + TRANSLATION_CHUNK_CHARS] for i in range(0, len(text), TRANSLATION_CHUNK_CHARS)]
    translated = [_translate_chunk(chunk) for chunk in chunks]
    return "\n\n".join(translated).strip()


def _normalize(value: str) -> str:
    decomposed = unicodedata.normalize("NFKD", value or "")
    ascii_value = "".join(c for c in decomposed if not unicodedata.combining(c))
    return re.sub(r"[^a-z0-9]+", " ", ascii_value.casefold()).strip()


def _normalized_reference(value: str | None) -> str:
    return str(value or "").strip().casefold().rstrip("/")


def find_duplicate(db: Session, analysis: dict, *, fingerprint: str | None = None) -> Document | None:
    """Deduplica por DOI, URL/fonte, metadado de título e fingerprint do arquivo.

    O fingerprint é aplicado novamente imediatamente antes da incorporação para
    impedir que duas bibliotecas privadas publiquem cópias do mesmo arquivo com
    metadados ligeiramente diferentes.
    """
    doi = _normalized_reference(str(analysis.get("doi") or ""))
    source_url = _normalized_reference(str(analysis.get("source_url") or ""))
    fingerprint_ref = f"sha256:{str(fingerprint or '').strip().casefold()}" if fingerprint else ""
    title = _normalize(str(analysis.get("title") or ""))

    for document in db.query(Document).all():
        refs = {_normalized_reference(str(ref)) for ref in (document.source_refs or []) if str(ref).strip()}
        if doi and any(
            ref == doi
            or ref == f"doi:{doi}"
            or ref == f"https://doi.org/{doi}"
            or ref == f"http://doi.org/{doi}"
            for ref in refs
        ):
            return document
        if source_url and source_url in refs:
            return document
        if fingerprint_ref and fingerprint_ref in refs:
            return document
        if title and _normalize(document.title) == title:
            return document
    return None


def _slug(value: str) -> str:
    normalized = _normalize(value).replace(" ", "-")[:180].strip("-") or "documento-cientifico"
    return normalized


def _kind(document_type: str) -> str:
    return {
        "diretriz": "diretriz",
        "consenso": "consenso",
        "ensaio_clinico": "estudo",
        "estudo_observacional": "estudo",
        "revisao_sistematica": "estudo",
        "meta_analise": "estudo",
        "evidencia": "estudo",
        "artigo": "estudo",
        "outro": "modulo",
    }.get(document_type, "modulo")


def has_traceable_source(analysis: dict) -> bool:
    return bool(str(analysis.get("doi") or "").strip() or str(analysis.get("source_url") or "").strip())


def _shared_synthesis_body(analysis: dict) -> str:
    """Conteúdo original do CorVIA para o acervo compartilhado.

    A tradução integral do arquivo enviado nunca sai da biblioteca privada do
    assinante. O corpus global recebe somente síntese derivada e rastreável.
    """
    parts = [
        "## Síntese clínica CorVIA",
        str(analysis.get("summary_pt") or "").strip(),
        "\n## Metodologia",
        str(analysis.get("methodology_pt") or "").strip(),
        "\n## População",
        str(analysis.get("population_pt") or "").strip(),
        "\n## Intervenções / exposições",
        str(analysis.get("interventions_pt") or "").strip(),
        "\n## Desfechos e resultados",
        "\n\n".join(filter(None, [
            str(analysis.get("outcomes_pt") or "").strip(),
            str(analysis.get("results_pt") or "").strip(),
        ])),
        "\n## Pontos-chave",
        "\n".join(f"- {x}" for x in analysis.get("key_points_pt") or []),
        "\n## Implicações clínicas",
        "\n".join(f"- {x}" for x in analysis.get("clinical_implications_pt") or []),
        "\n## Limitações",
        "\n".join(f"- {x}" for x in analysis.get("limitations_pt") or []),
        "\n## Proveniência",
        "Síntese original produzida pelo CorVIA a partir de documento enviado por assinante e incorporado mediante consentimento explícito. Consulte a fonte primária para o texto integral.",
    ]
    return "\n\n".join(part for part in parts if part.strip()).strip()


def incorporate(db: Session, row: ScientificUserDocument, analysis: dict, translated_text: str, *, reviewer_id: int) -> Document:
    duplicate = find_duplicate(db, analysis, fingerprint=row.sha256)
    if duplicate:
        row.incorporation_status = "duplicado"
        row.incorporated_document_id = duplicate.id
        return duplicate
    if not has_traceable_source(analysis):
        raise ValueError("Documento sem DOI ou URL de fonte rastreável não pode ser incorporado automaticamente ao acervo compartilhado.")

    base_slug = _slug(str(analysis.get("title") or "Documento científico"))
    slug = base_slug
    suffix = 2
    while db.query(Document.id).filter(Document.slug == slug).first():
        slug = f"{base_slug[:170]}-{suffix}"
        suffix += 1

    topics = [str(x).strip() for x in analysis.get("topics") or [] if str(x).strip()]
    theme = (topics[0] if topics else "Geral")[:80]
    refs = []
    for value in (analysis.get("doi"), analysis.get("source_url")):
        value = str(value or "").strip()
        if value and value not in refs:
            refs.append(value)
    refs.append(f"sha256:{row.sha256}")

    document = Document(
        slug=slug,
        title=str(analysis.get("title") or "Documento científico")[:500],
        kind=_kind(str(analysis.get("document_type") or row.document_type)),
        theme=theme,
        summary=str(analysis.get("summary_pt") or ""),
        body_md=_shared_synthesis_body(analysis),
        tags=topics[:30],
        source_refs=refs,
        evidence_level=str(analysis.get("evidence_strength") or "nao_determinada")[:40],
        source_tier="A",
        review_status="revisado",
        published=True,
        reviewed_by=reviewer_id,
        reviewed_at=datetime.now(timezone.utc),
        gaps=[],
    )
    db.add(document)
    db.flush()
    row.incorporated_document_id = document.id
    row.incorporation_status = "incorporado"
    backfill_mesmo_tema(db, commit=False)
    return document


def decrypt_payload(row: ScientificUserDocument) -> tuple[str, str, dict]:
    extracted = cofre.decifrar_campo(row.extracted_text_cifrado, row.id) if row.extracted_text_cifrado else ""
    translated = cofre.decifrar_campo(row.translated_text_cifrado, row.id) if row.translated_text_cifrado else ""
    if row.analysis_cifrado:
        analysis = json.loads(cofre.decifrar_campo(row.analysis_cifrado, row.id))
    else:
        analysis = {}
    return extracted, translated, analysis
