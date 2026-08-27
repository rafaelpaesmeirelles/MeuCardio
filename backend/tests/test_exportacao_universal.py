"""Exportação universal: conteúdo canônico → PDF/PPTX/DOCX ou CorVIA Mail.

Os testes usam um documento clínico mínimo publicado e não dependem do corpus
carregado no banco de teste. A rota é exercitada de ponta a ponta com a mesma
dependência global de assinatura principal usada em produção.
"""
import datetime
from io import BytesIO

import pytest
from cryptography import x509
from cryptography.hazmat.primitives import hashes, serialization
from cryptography.hazmat.primitives.asymmetric import rsa
from cryptography.hazmat.primitives.serialization import pkcs12
from cryptography.x509.oid import NameOID
from docx import Document as WordDocument
from pptx import Presentation
from sqlalchemy import text

from app.core.config import settings
from app.models.content import Document
from app.models.email_account import EmailAccount
from app.models.subscription import Subscription
from app.services.assinatura import certificado_a1


PREFIXO = "teste-exportacao-universal-"


@pytest.fixture(autouse=True)
def _ambiente_exportacao(db, tmp_path, monkeypatch):
    monkeypatch.setattr(settings, "certificados_dir", str(tmp_path / "certificados"))
    db.execute(text("DELETE FROM documents WHERE slug LIKE :prefixo"), {"prefixo": f"{PREFIXO}%"})
    db.commit()
    yield
    db.execute(text("DELETE FROM documents WHERE slug LIKE :prefixo"), {"prefixo": f"{PREFIXO}%"})
    db.commit()


def _headers(token: str) -> dict[str, str]:
    return {"Authorization": f"Bearer {token}"}


def _assinatura_principal(db, user):
    db.add(Subscription(user_id=user.id, kind="meucardio", plano="basico", status="ativo"))
    db.commit()


def _ativar_corvia_mail(db, user, endereco: str):
    db.add(Subscription(user_id=user.id, kind="email", status="ativo"))
    db.add(EmailAccount(
        user_id=user.id,
        email_address=endereco,
        mail360_account_key=f"mail-key-{user.id}",
        status="ativa",
    ))
    db.commit()


def _gerar_pfx(*, senha: str = "senha123") -> bytes:
    chave = rsa.generate_private_key(public_exponent=65537, key_size=2048)
    nome = x509.Name([x509.NameAttribute(NameOID.COMMON_NAME, "DR EXPORTADOR:12345678900")])
    agora = datetime.datetime.now(datetime.timezone.utc)
    certificado = (
        x509.CertificateBuilder()
        .subject_name(nome)
        .issuer_name(nome)
        .public_key(chave.public_key())
        .serial_number(x509.random_serial_number())
        .not_valid_before(agora - datetime.timedelta(days=1))
        .not_valid_after(agora + datetime.timedelta(days=365))
        .sign(chave, hashes.SHA256())
    )
    return pkcs12.serialize_key_and_certificates(
        name=b"exportacao-teste",
        key=chave,
        cert=certificado,
        cas=None,
        encryption_algorithm=serialization.BestAvailableEncryption(senha.encode()),
    )


def _documento(db, sufixo: str, *, publicado: bool = True) -> Document:
    doc = Document(
        slug=f"{PREFIXO}{sufixo}",
        title="Insuficiência cardíaca — documento de teste",
        kind="modulo",
        theme="Insuficiência cardíaca",
        summary="Resumo clínico validado para o teste de exportação.",
        body_md="## Conduta\n\nTexto clínico publicado.\n\n## Pontos-chave\n\n- Primeiro ponto\n- Segundo ponto",
        source_refs=["Diretriz de teste 2026"],
        review_status="revisado",
        published=publicado,
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)
    return doc


class TestCatalogoEPdf:
    def test_catalogo_encontra_item_e_pdf_e_real(self, client, db, criar_usuario):
        user, token = criar_usuario(email="exportador@teste.local")
        _assinatura_principal(db, user)
        doc = _documento(db, "pdf")

        catalogo = client.get(
            f"/api/exportar/catalogo?tipo=documento&slug={doc.slug}&limite=1",
            headers=_headers(token),
        )
        assert catalogo.status_code == 200, catalogo.text
        assert catalogo.json()["itens"][0]["slug"] == doc.slug

        resposta = client.post(
            "/api/exportar/conteudo",
            headers=_headers(token),
            json={
                "itens": [{"tipo": "documento", "slug": doc.slug}],
                "incluir_dados_assinante": True,
            },
        )
        assert resposta.status_code == 200, resposta.text
        assert resposta.headers["content-type"].startswith("application/pdf")
        assert resposta.content.startswith(b"%PDF")
        assert len(resposta.content) > 1000

    @pytest.mark.parametrize(
        ("formato", "content_type"),
        [
            ("pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
            ("docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
        ],
    )
    def test_exporta_formato_office_nativo_editavel(
        self, client, db, criar_usuario, formato, content_type,
    ):
        user, token = criar_usuario(email=f"exportador-{formato}@teste.local")
        _assinatura_principal(db, user)
        doc = _documento(db, formato)

        resposta = client.post(
            "/api/exportar/conteudo",
            headers=_headers(token),
            json={
                "itens": [{"tipo": "documento", "slug": doc.slug}],
                "incluir_dados_assinante": True,
                "formato": formato,
            },
        )

        assert resposta.status_code == 200, resposta.text
        assert resposta.headers["content-type"] == content_type
        assert resposta.headers["content-disposition"].endswith(f'.{formato}"')
        assert resposta.content[:2] == b"PK"

        if formato == "pptx":
            arquivo = Presentation(BytesIO(resposta.content))
            texto = "\n".join(
                shape.text_frame.text
                for slide in arquivo.slides
                for shape in slide.shapes
                if shape.has_text_frame
            )
            assert len(arquivo.slides) >= 5
        else:
            arquivo = WordDocument(BytesIO(resposta.content))
            texto = "\n".join(paragrafo.text for paragrafo in arquivo.paragraphs)
            assert arquivo.sections

        assert "Insuficiência cardíaca" in texto
        assert "Texto clínico publicado" in texto
        assert "Diretriz de teste 2026" in texto
        assert "não cria nem atualiza recomendações clínicas" in texto

    def test_formato_invalido_e_rejeitado(self, client, db, criar_usuario):
        user, token = criar_usuario(email="exportador-invalido@teste.local")
        _assinatura_principal(db, user)
        doc = _documento(db, "invalido")

        resposta = client.post(
            "/api/exportar/conteudo",
            headers=_headers(token),
            json={
                "itens": [{"tipo": "documento", "slug": doc.slug}],
                "formato": "txt",
            },
        )
        assert resposta.status_code == 422

    def test_pdf_assinado_cobre_arquivo_e_exibe_selo_em_todas_as_paginas(
        self, client, db, criar_usuario,
    ):
        import fitz
        from pyhanko.pdf_utils.reader import PdfFileReader
        from pyhanko.sign.validation import validate_pdf_signature

        user, token = criar_usuario(email="exportador-assinado@teste.local")
        _assinatura_principal(db, user)
        certificado_a1.salvar(db, user, _gerar_pfx(), "senha123", None)
        db.commit()
        documentos = [_documento(db, f"assinado-{indice}") for indice in range(1, 7)]

        resposta = client.post(
            "/api/exportar/conteudo",
            headers=_headers(token),
            json={
                "itens": [{"tipo": "documento", "slug": doc.slug} for doc in documentos],
                "incluir_dados_assinante": True,
                "formato": "pdf",
                "assinar_digitalmente": True,
            },
        )

        assert resposta.status_code == 200, resposta.text
        leitor = PdfFileReader(BytesIO(resposta.content))
        assert len(leitor.embedded_signatures) == 1
        status = validate_pdf_signature(leitor.embedded_signatures[0])
        assert status.intact is True
        assert status.valid is True

        documento = fitz.open(stream=resposta.content, filetype="pdf")
        assert len(documento) >= 2
        assert documento[0].first_widget is not None
        for pagina in documento[1:]:
            assert any(nome.startswith("Stamp") for _xref, nome, *_resto in pagina.get_xobjects())
            assert "ASSINATURA DIGITAL ICP-BRASIL" in pagina.get_text()

    def test_assinatura_pdf_falha_fechada_sem_certificado(self, client, db, criar_usuario):
        user, token = criar_usuario(email="exportador-sem-certificado@teste.local")
        _assinatura_principal(db, user)
        doc = _documento(db, "sem-certificado")

        resposta = client.post(
            "/api/exportar/conteudo",
            headers=_headers(token),
            json={
                "itens": [{"tipo": "documento", "slug": doc.slug}],
                "formato": "pdf",
                "assinar_digitalmente": True,
            },
        )

        assert resposta.status_code == 409, resposta.text
        assert "certificado A1" in resposta.text

    @pytest.mark.parametrize("formato", ["pptx", "docx"])
    def test_nao_simula_assinatura_em_formato_office(
        self, client, db, criar_usuario, formato,
    ):
        user, token = criar_usuario(email=f"exportador-assinatura-{formato}@teste.local")
        _assinatura_principal(db, user)
        doc = _documento(db, f"assinatura-{formato}")

        resposta = client.post(
            "/api/exportar/conteudo",
            headers=_headers(token),
            json={
                "itens": [{"tipo": "documento", "slug": doc.slug}],
                "formato": formato,
                "assinar_digitalmente": True,
            },
        )

        assert resposta.status_code == 422, resposta.text
        assert "somente para PDF" in resposta.text

    def test_fail_closed_quando_item_nao_esta_publicado(self, client, db, criar_usuario):
        user, token = criar_usuario(email="exportador-retido@teste.local")
        _assinatura_principal(db, user)
        doc = _documento(db, "retido", publicado=False)

        resposta = client.post(
            "/api/exportar/conteudo",
            headers=_headers(token),
            json={"itens": [{"tipo": "documento", "slug": doc.slug}]},
        )
        assert resposta.status_code == 409, resposta.text
        assert doc.slug in resposta.text


class TestCorviaMail:
    def test_status_indisponivel_sem_addon(self, client, db, criar_usuario):
        user, token = criar_usuario(email="sem-mail-export@teste.local")
        _assinatura_principal(db, user)

        resposta = client.get("/api/exportar/corvia-mail", headers=_headers(token))
        assert resposta.status_code == 200, resposta.text
        assert resposta.json()["disponivel"] is False

    def test_caixa_nativa_falha_fechada_quando_smime_esta_ativo(
        self, client, db, criar_usuario, monkeypatch_mail360,
    ):
        user, token = criar_usuario(email="smime-export@teste.local")
        _assinatura_principal(db, user)
        _ativar_corvia_mail(db, user, "smime-export@corvia.med.br")
        doc = _documento(db, "smime-bloqueado")
        user.email_assinatura_digital_ativa = True
        db.commit()

        status = client.get("/api/exportar/corvia-mail", headers=_headers(token))
        assert status.status_code == 200, status.text
        assert status.json()["disponivel"] is False
        assert "S/MIME" in status.json()["motivo"]

        resposta = client.post(
            "/api/exportar/conteudo/enviar-email",
            headers=_headers(token),
            json={
                "itens": [{"tipo": "documento", "slug": doc.slug}],
                "para": "destinatario@teste.local",
            },
        )
        assert resposta.status_code == 409, resposta.text
        assert "S/MIME" in resposta.text
        assert monkeypatch_mail360["anexos"] == []
        assert monkeypatch_mail360["mensagens_enviadas"] == []

    def test_gera_anexa_e_envia_pela_caixa_do_assinante(
        self, client, db, criar_usuario, monkeypatch_mail360,
    ):
        user, token = criar_usuario(email="com-mail-export@teste.local")
        _assinatura_principal(db, user)
        _ativar_corvia_mail(db, user, "medico-export@corvia.med.br")
        doc = _documento(db, "email")

        status = client.get("/api/exportar/corvia-mail", headers=_headers(token))
        assert status.status_code == 200, status.text
        assert status.json() == {
            "disponivel": True,
            "motivo": None,
            "email_address": "medico-export@corvia.med.br",
        }

        resposta = client.post(
            "/api/exportar/conteudo/enviar-email",
            headers=_headers(token),
            json={
                "itens": [{"tipo": "documento", "slug": doc.slug}],
                "incluir_dados_assinante": True,
                "para": "destinatario@teste.local",
                "cc": "copia@teste.local",
                "assunto": "Conteúdo para discussão",
                "mensagem": "Segue o PDF selecionado no CorVIA.",
            },
        )
        assert resposta.status_code == 201, resposta.text
        corpo = resposta.json()
        assert corpo["enviado"] is True
        assert corpo["remetente"] == "medico-export@corvia.med.br"
        assert corpo["para"] == "destinatario@teste.local"
        assert corpo["message_id"] == "msg-enviada-1"

        assert len(monkeypatch_mail360["anexos"]) == 1
        account_key, nome, tamanho = monkeypatch_mail360["anexos"][0]
        assert account_key == f"mail-key-{user.id}"
        assert nome.endswith(".pdf")
        assert tamanho > 1000

        assert len(monkeypatch_mail360["mensagens_enviadas"]) == 1
        enviada = monkeypatch_mail360["mensagens_enviadas"][0]
        assert enviada["remetente"] == "medico-export@corvia.med.br"
        assert enviada["para"] == "destinatario@teste.local"
        assert enviada["cc"] == "copia@teste.local"
        assert enviada["anexos"] == ["file-id-1"]

    def test_destinatario_invalido_e_rejeitado_antes_do_envio(
        self, client, db, criar_usuario, monkeypatch_mail360,
    ):
        user, token = criar_usuario(email="mail-invalido-export@teste.local")
        _assinatura_principal(db, user)
        _ativar_corvia_mail(db, user, "medico-invalid@corvia.med.br")
        doc = _documento(db, "email-invalido")

        resposta = client.post(
            "/api/exportar/conteudo/enviar-email",
            headers=_headers(token),
            json={
                "itens": [{"tipo": "documento", "slug": doc.slug}],
                "para": "endereco-invalido",
            },
        )
        assert resposta.status_code == 422, resposta.text
        assert monkeypatch_mail360["anexos"] == []
        assert monkeypatch_mail360["mensagens_enviadas"] == []
