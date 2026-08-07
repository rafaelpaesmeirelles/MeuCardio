"""Modo apresentação em PowerPoint editável (pedido do Rafael, 07/08/2026: dar
a opção de PDF ou .pptx editável no Modo Apresentação). Testado pela rota HTTP
real — mesmo conteúdo do PDF, formato de arquivo diferente.
"""
import zipfile
from io import BytesIO

from pptx import Presentation

from app.models.content import Document
from app.models.subscription import Subscription


def _headers(token: str) -> dict[str, str]:
    return {"Authorization": f"Bearer {token}"}


def _subscribe(db, user_id: int) -> None:
    db.add(Subscription(user_id=user_id, kind="meucardio", plano="basico", status="ativo"))
    db.commit()


def _publicar_documento(db) -> Document:
    doc = Document(
        slug="documento-de-teste-apresentacao",
        title="Documento de teste para apresentação",
        theme="Geral",
        summary="Resumo de teste.",
        body_md=(
            "## Primeira seção\n\n"
            "- item um\n"
            "- item dois\n\n"
            "Parágrafo de corpo com alguma frase mais longa para exercitar a "
            "fragmentação de texto do exportador.\n\n"
            "## Segunda seção\n\n"
            "- item três\n"
        ),
        source_refs=["Fonte de teste. Revista. 2024;1(1):1-2."],
        review_status="revisado",
        published=True,
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)
    return doc


def test_apresentacao_formato_padrao_continua_pdf(client, db, criar_usuario):
    _publicar_documento(db)
    user, token = criar_usuario()
    _subscribe(db, user.id)
    resposta = client.post(
        "/api/biblioteca/documento-de-teste-apresentacao/apresentacao",
        json={"anotacao": ""},
        headers=_headers(token),
    )
    assert resposta.status_code == 200
    assert resposta.headers["content-type"] == "application/pdf"
    assert resposta.content.startswith(b"%PDF-")


def test_apresentacao_formato_pptx_gera_arquivo_powerpoint_valido(client, db, criar_usuario):
    _publicar_documento(db)
    user, token = criar_usuario()
    _subscribe(db, user.id)
    resposta = client.post(
        "/api/biblioteca/documento-de-teste-apresentacao/apresentacao",
        json={"anotacao": "", "formato": "pptx"},
        headers=_headers(token),
    )
    assert resposta.status_code == 200
    assert resposta.headers["content-type"] == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    # .pptx é um ZIP OOXML — confere a assinatura de arquivo real, não só o header.
    assert resposta.content[:2] == b"PK"
    with zipfile.ZipFile(BytesIO(resposta.content)) as z:
        assert "ppt/presentation.xml" in z.namelist()

    # Abre de verdade com python-pptx e confere que o conteúdo do documento
    # (não texto genérico) chegou aos slides.
    prs = Presentation(BytesIO(resposta.content))
    assert len(prs.slides) >= 4  # capa + 2 seções + procedência, no mínimo
    texto_completo = "\n".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Documento de teste para apresentação" in texto_completo
    assert "item um" in texto_completo
    assert "item três" in texto_completo
    assert "Fonte de teste" in texto_completo


def test_apresentacao_formato_invalido_devolve_422(client, db, criar_usuario):
    _publicar_documento(db)
    user, token = criar_usuario()
    _subscribe(db, user.id)
    resposta = client.post(
        "/api/biblioteca/documento-de-teste-apresentacao/apresentacao",
        json={"formato": "docx"},
        headers=_headers(token),
    )
    assert resposta.status_code == 422


def test_apresentacao_pptx_com_anotacao_inclui_slide_de_anotacao(client, db, criar_usuario):
    _publicar_documento(db)
    user, token = criar_usuario()
    _subscribe(db, user.id)
    resposta = client.post(
        "/api/biblioteca/documento-de-teste-apresentacao/apresentacao",
        json={"anotacao": "Discutir o caso do leito 12.", "formato": "pptx"},
        headers=_headers(token),
    )
    assert resposta.status_code == 200
    prs = Presentation(BytesIO(resposta.content))
    texto_completo = "\n".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Discutir o caso do leito 12" in texto_completo
    assert "não faz parte do conteúdo verificado" in texto_completo
